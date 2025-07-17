# pdf_tools/pdf_to_pptx_logic.py
import io
import os
import fitz # PyMuPDF
import re
import tempfile
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

DPI_OPTIONS = { "72": 72, "150": 150, "300": 300 }

def parse_page_ranges(page_string, total_pages):
    # (Same function as before)
    indices = set()
    if not page_string: return None, "Page range string cannot be empty."
    parts = re.split(r'[,\s]+', page_string)
    for part in parts:
        part = part.strip();
        if not part: continue
        if '-' in part:
            try:
                start, end = map(int, part.split('-', 1))
                if start < 1 or end < 1: return None, f"Page numbers must be positive: '{part}'."
                if start > end: return None, f"Invalid range: start > end in '{part}'."
                if start > total_pages or end > total_pages: return None, f"Page number out of bounds (max {total_pages}): '{part}'."
                indices.update(range(start - 1, end))
            except ValueError: return None, f"Invalid range format: '{part}'. Use '3-5'."
        else:
            try:
                page_num = int(part)
                if page_num < 1: return None, f"Page number must be positive: '{part}'."
                if page_num > total_pages: return None, f"Page number out of bounds (max {total_pages}): '{part}'."
                indices.add(page_num - 1)
            except ValueError: return None, f"Invalid page number: '{part}'. Use numbers or '1, 3-5'."
    if not indices: return None, "No valid pages selected."
    return sorted(list(indices)), None


def convert_pdf_to_pptx_images(pdf_file, dpi_str='150', page_selection_mode='all', page_string=None):
    if not pdf_file:
        return None, None, "No PDF file provided."

    dpi = DPI_OPTIONS.get(dpi_str, 150)
    output_buffer = io.BytesIO()
    output_filename = "converted_presentation.pptx"

    temp_pdf_path = None # Variable for temporary file path
    prs = None

    try:
        # --- Use Temporary File ---
        pdf_file.seek(0)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        print(f"[PDF->PPTX Logic] Saved temp PDF to: {temp_pdf_path}")

        # --- Get Page Count and Check Encryption (Open Once) ---
        doc_check = None
        try:
            doc_check = fitz.open(temp_pdf_path)
            total_pages = len(doc_check)
            if doc_check.is_encrypted:
                return None, None, f"File '{pdf_file.name}' is encrypted and password protected."
        except Exception as e:
             print(f"[PDF->PPTX Logic] Error opening PDF for check: {e}")
             return None, None, f"Could not read PDF file: {e}"
        finally:
            if doc_check: doc_check.close()
        # --- End Check ---

        print(f"[PDF->PPTX Logic] Total pages: {total_pages}")

        pages_to_convert = []
        if page_selection_mode == 'specific':
            parsed_indices, error = parse_page_ranges(page_string, total_pages)
            if error: return None, None, error
            pages_to_convert = parsed_indices
            print(f"[PDF->PPTX Logic] Specific pages selected (0-based): {pages_to_convert}")
        else:
            pages_to_convert = list(range(total_pages))
            print(f"[PDF->PPTX Logic] Converting all {total_pages} pages.")

        if not pages_to_convert:
             return None, None, "No pages selected or found for conversion."

        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        pages_added_count = 0

        # --- Read the PDF bytes once ---
        # This avoids potential issues with the temp file being accessed concurrently
        with open(temp_pdf_path, "rb") as f:
             pdf_data_bytes = f.read()

        for i in pages_to_convert:
            page_doc = None # Document object for this page only
            img_buffer = None # Buffer for this page only
            pil_img = None # PIL image for this page only
            try:
                # --- Open PDF stream *for each page* --- (Inefficient, but trying to isolate issue)
                page_doc = fitz.open(stream=pdf_data_bytes, filetype="pdf")
                if i >= len(page_doc): # Sanity check index
                    print(f"[PDF->PPTX Logic] Page index {i} out of bounds after re-opening doc. Skipping.")
                    continue

                page = page_doc.load_page(i)
                print(f"[PDF->PPTX Logic] Rendering page {i+1} at {dpi} DPI...")
                pix = page.get_pixmap(matrix=fitz.Matrix(dpi/72, dpi/72), alpha=False)

                if pix.samples:
                    # --- Convert pixmap to PIL Image and save to buffer ---
                    img_buffer = io.BytesIO()
                    # Use RGB for JPG saving with Pillow
                    pil_img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    pil_img.save(img_buffer, format="PNG") # Save as PNG intermediate
                    img_buffer.seek(0)
                    # --- End PIL Conversion ---

                    slide = prs.slides.add_slide(blank_slide_layout)
                    # --- Add picture logic (same as before) ---
                    slide_width_in = Inches(prs.slide_width / 914400)
                    slide_height_in = Inches(prs.slide_height / 914400)
                    img_width_px, img_height_px = pil_img.size # Get dimensions from PIL image
                    img_aspect = img_width_px / img_height_px if img_height_px > 0 else 1
                    slide_aspect = slide_width_in / slide_height_in if slide_height_in > 0 else 1
                    if img_aspect > slide_aspect: width_in=slide_width_in; height_in=width_in/img_aspect
                    else: height_in=slide_height_in; width_in=height_in*img_aspect
                    left_in = (slide_width_in - width_in) / 2
                    top_in = (slide_height_in - height_in) / 2
                    slide.shapes.add_picture(img_buffer, left_in, top_in, width=width_in, height=height_in)
                    # --- End Add Picture Logic ---

                    print(f"[PDF->PPTX Logic] Added image for page {i+1} to slide.")
                    pages_added_count += 1
                else:
                    print(f"[PDF->PPTX Logic] Warning: Failed to get pixel data for page {i+1}.")

            except Exception as page_err:
                 print(f"[PDF->PPTX Logic] Error processing page {i+1}: {page_err}")
                 # Continue to next page
            finally:
                 # Clean up resources FOR THIS PAGE's loop iteration
                 if pil_img:
                     try: pil_img.close()
                     except Exception: pass
                 if img_buffer:
                     try: img_buffer.close()
                     except Exception: pass
                 if page_doc: # Close the doc opened for this page
                      try: page_doc.close()
                      except Exception: pass


        if pages_added_count == 0:
             return None, None, "No pages could be successfully converted."

        # Save presentation
        prs.save(output_buffer)
        output_buffer.seek(0)

        base_name = os.path.splitext(pdf_file.name)[0]
        output_filename = f"{base_name}_presentation.pptx"

        print("[PDF->PPTX Logic] Presentation created successfully.")
        return output_buffer.getvalue(), output_filename, None

    except Exception as e:
        print(f"Error during PDF to PPTX conversion process: {e}")
        import traceback
        traceback.print_exc()
        return None, None, f"An unexpected error occurred: {e}"
    finally:
        # Clean up temporary file
        if temp_pdf_path and os.path.exists(temp_pdf_path):
            try: os.remove(temp_pdf_path); print(f"[PDF->PPTX Cleanup] Deleted temp PDF: {temp_pdf_path}")
            except OSError as e: print(f"[PDF->PPTX Cleanup] Error deleting temp PDF {temp_pdf_path}: {e}")